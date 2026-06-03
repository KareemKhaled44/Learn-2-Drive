from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Learn_2_Drive_Graduation_Presentation.pptx"

BG = "0F172A"
PANEL = "111827"
PANEL_2 = "1E293B"
CYAN = "22D3EE"
BLUE = "1E40AF"
SLATE = "94A3B8"
WHITE = "F8FAFC"
MUTED = "CBD5E1"
GOOD = "34D399"
GOLD = "F59E0B"
RED = "F87171"


def rgb(hex_value: str):
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_background(slide, title_band=False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    width = prs.slide_width
    height = prs.slide_height

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, width, Inches(0.24))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(CYAN)
    bar.line.fill.background()

    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.25), Inches(0.35), Inches(0.08), Inches(6.65))
    side.fill.solid()
    side.fill.fore_color.rgb = rgb(BLUE)
    side.line.fill.background()

    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.9), Inches(-0.4), Inches(3.0), Inches(3.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = rgb(BLUE)
    glow.fill.transparency = 0.82
    glow.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.6), Inches(5.3), Inches(2.4), Inches(2.4))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = rgb(CYAN)
    glow2.fill.transparency = 0.86
    glow2.line.fill.background()

    if title_band:
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.15), Inches(5.6), Inches(0.6))
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(PANEL_2)
        band.line.color.rgb = rgb(CYAN)
        band.line.width = Pt(1.5)


def add_footer(slide, index):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(7.0), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(CYAN)
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.62), Inches(7.04), Inches(8.0), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Learn 2 Drive | BIS Graduation Project"
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.color.rgb = rgb(SLATE)

    page = slide.shapes.add_textbox(Inches(12.15), Inches(7.02), Inches(0.6), Inches(0.25))
    tf = page.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(index)
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = rgb(CYAN)


def add_heading(slide, title, subtitle=None, top=0.52, size=24):
    box = slide.shapes.add_textbox(Inches(0.62), Inches(top), Inches(8.8), Inches(0.85))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = rgb(WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Aptos"
        p2.font.size = Pt(11)
        p2.font.color.rgb = rgb(MUTED)
        p2.space_before = Pt(4)


def add_panel(slide, left, top, width, height, fill=PANEL, line=BLUE, radius=True):
    shape = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(fill)
    box.line.color.rgb = rgb(line)
    box.line.width = Pt(1.25)
    return box


def add_text(slide, left, top, width, height, text, font_size=16, color=WHITE, bold=False, font="Aptos", align=PP_ALIGN.LEFT, margin=0.08):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(font_size)
    p.font.color.rgb = rgb(color)
    p.font.bold = bold
    p.alignment = align
    return box


def add_multiline_text(slide, left, top, width, height, lines, font_size=15, color=WHITE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb(color)
        p.line_spacing = Pt(font_size * 1.05)
    return box


def add_metric_card(slide, left, top, width, height, label, value, caption=None, accent=CYAN):
    add_panel(slide, left, top, width, height, fill=PANEL_2, line=accent)
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.25), label, font_size=10, color=MUTED)
    add_text(slide, left + Inches(0.18), top + Inches(0.36), width - Inches(0.36), Inches(0.4), value, font_size=22, color=WHITE, bold=True)
    if caption:
        add_text(slide, left + Inches(0.18), top + Inches(0.9), width - Inches(0.36), Inches(0.4), caption, font_size=10, color=SLATE)


def add_chip(slide, left, top, width, text, fill=BLUE, font_size=10):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.32))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(fill)
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = rgb(WHITE)
    return chip


def add_role_card(slide, left, top, width, height, title, body, accent, symbol):
    add_panel(slide, left, top, width, height, fill=PANEL_2, line=accent)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), top + Inches(0.18), Inches(0.52), Inches(0.52))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(accent)
    badge.line.fill.background()
    add_text(slide, left + Inches(0.18), top + Inches(0.20), Inches(0.52), Inches(0.28), symbol, font_size=18, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.84), top + Inches(0.16), width - Inches(1.0), Inches(0.25), title, font_size=15, bold=True)
    add_multiline_text(slide, left + Inches(0.18), top + Inches(0.68), width - Inches(0.36), height - Inches(0.8), body, font_size=11, color=MUTED)


def add_transition(slide):
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    slide_element = slide._element
    transition = OxmlElement("p:transition")
    transition.set("advClick", "1")
    transition.set("spd", "med")
    transition.append(OxmlElement("p:fade"))
    children = list(slide_element)
    insert_at = len(children)
    for i, child in enumerate(children):
        if child.tag.endswith("clrMapOvr"):
            insert_at = i
            break
    slide_element.insert(insert_at, transition)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title_band=True)

    add_chip(slide, Inches(0.85), Inches(0.48), Inches(1.5), "GRADUATION DEFENSE", fill=BLUE)
    add_text(slide, Inches(0.82), Inches(1.38), Inches(7.6), Inches(1.0), "Learn 2 Drive", font_size=34, bold=True)
    add_text(slide, Inches(0.84), Inches(2.28), Inches(7.8), Inches(0.9), "A digital platform for driving academy discovery and booking", font_size=16, color=MUTED)

    add_metric_card(slide, Inches(0.86), Inches(3.05), Inches(2.0), Inches(1.15), "Project Year", "2025–2026", "BIS Graduation Project")
    add_metric_card(slide, Inches(3.08), Inches(3.05), Inches(2.0), Inches(1.15), "Team Members", "1. 2. 3. 4. 5.", "Insert final names")
    add_metric_card(slide, Inches(5.30), Inches(3.05), Inches(2.2), Inches(1.15), "Supervisors", "IT + Business", "Dr. ___________________")

    add_panel(slide, Inches(8.25), Inches(1.18), Inches(4.15), Inches(4.8), fill=PANEL, line=CYAN)
    add_text(slide, Inches(8.52), Inches(1.42), Inches(3.5), Inches(0.35), "Why this project matters", font_size=16, bold=True)
    add_multiline_text(slide, Inches(8.55), Inches(1.84), Inches(3.45), Inches(2.15), [
        "• Centralizes academy discovery in one place",
        "• Replaces manual booking with online flows",
        "• Adds approval, ratings, and transparent comparison",
        "• Supports guest browsing and role-based dashboards",
    ], font_size=13, color=MUTED)

    roadmap = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.52), Inches(4.25), Inches(3.45), Inches(1.15))
    roadmap.fill.solid()
    roadmap.fill.fore_color.rgb = rgb(BLUE)
    roadmap.line.fill.background()
    add_text(slide, Inches(8.7), Inches(4.44), Inches(3.1), Inches(0.25), "Modern defense-ready system design", font_size=14, bold=True)
    add_text(slide, Inches(8.7), Inches(4.72), Inches(3.0), Inches(0.20), "Dark cyan palette | clean visuals | clear flow", font_size=10, color=WHITE)

    add_footer(slide, 1)
    add_transition(slide)


def add_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Agenda", "15 content slides covering the project from problem to future work")

    items = [
        ("01", "Title", "Project identity and defense context"),
        ("02", "Agenda", "Roadmap of the presentation"),
        ("03", "Business Problem", "Current market pain points"),
        ("04", "Project Idea & Scope", "What the platform solves"),
        ("05", "Related Works & Gap", "Why Learn 2 Drive is different"),
        ("06", "System Features & Users", "Actors and responsibilities"),
        ("07", "Business Plan Highlights", "Go-to-market and value model"),
        ("08", "Financial Snapshot", "Startup capital and Year 1 outlook"),
        ("09", "System Analysis", "Use case and activity flows"),
        ("10", "System Design", "ERD and database mapping"),
        ("11", "Architecture & Tools", "Technology stack and layers"),
        ("12", "Implementation Highlights", "Core screens and code parts"),
        ("13", "Key Technical Challenges", "Hard parts and solutions"),
        ("14", "Testing Results", "47/47 test cases passed"),
        ("15", "Conclusion & Future Work", "Summary and next steps"),
    ]

    for idx, (num, title, desc) in enumerate(items):
        row = idx // 3
        col = idx % 3
        left = Inches(0.82 + col * 4.12)
        top = Inches(1.35 + row * 2.0)
        add_panel(slide, left, top, Inches(3.76), Inches(1.55), fill=PANEL_2, line=BLUE)
        add_chip(slide, left + Inches(0.16), top + Inches(0.14), Inches(0.52), num, fill=CYAN, font_size=10)
        add_text(slide, left + Inches(0.84), top + Inches(0.12), Inches(2.6), Inches(0.25), title, font_size=15, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.60), Inches(3.25), Inches(0.55), desc, font_size=11, color=MUTED)

    add_footer(slide, 2)
    add_transition(slide)


def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Business Problem", "The report identifies a fragmented, manual market with limited transparency")

    add_panel(slide, Inches(0.82), Inches(1.52), Inches(6.05), Inches(4.9), fill=PANEL_2, line=CYAN)
    add_text(slide, Inches(1.05), Inches(1.72), Inches(2.2), Inches(0.25), "Core pain points", font_size=16, bold=True)
    bullets = [
        "Information is scattered across websites, social media, and word-of-mouth.",
        "Pricing, packages, and instructor quality are often not transparent.",
        "Booking is still manual and usually requires phone calls or in-person visits.",
        "Learners with specific needs, such as female instructors, cannot filter easily.",
        "There is no reliable review layer for comparing academy quality.",
    ]
    add_multiline_text(slide, Inches(1.05), Inches(2.08), Inches(5.45), Inches(3.55), [f"• {x}" for x in bullets], font_size=13, color=MUTED)

    add_panel(slide, Inches(7.1), Inches(1.52), Inches(5.35), Inches(4.9), fill=PANEL, line=BLUE)
    add_text(slide, Inches(7.35), Inches(1.72), Inches(2.6), Inches(0.25), "Business impact", font_size=16, bold=True)
    add_metric_card(slide, Inches(7.35), Inches(2.05), Inches(2.25), Inches(1.0), "Transparency", "Low", "Hidden pricing and service quality", accent=RED)
    add_metric_card(slide, Inches(9.75), Inches(2.05), Inches(2.25), Inches(1.0), "Booking", "Manual", "Offline, time-consuming process", accent=GOLD)
    add_metric_card(slide, Inches(7.35), Inches(3.20), Inches(2.25), Inches(1.0), "Search", "Fragmented", "No central comparison layer", accent=BLUE)
    add_metric_card(slide, Inches(9.75), Inches(3.20), Inches(2.25), Inches(1.0), "Feedback", "Missing", "No trust-building review system", accent=CYAN)

    add_text(slide, Inches(7.35), Inches(4.55), Inches(4.55), Inches(0.55), "The platform answers this gap with a centralized marketplace, online booking, filtering, ratings, and approval control.", font_size=13, color=MUTED)
    add_footer(slide, 3)
    add_transition(slide)


def add_scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Project Idea & Scope", "A centralized platform connecting learners, academies, trainers, and administrators")

    add_role_card(slide, Inches(0.84), Inches(1.6), Inches(3.0), Inches(1.5), "Central platform", [
        "Connects driving learners with academies in one online environment.",
        "Supports discovery, comparison, and booking.",
    ], CYAN, "C")
    add_role_card(slide, Inches(4.03), Inches(1.6), Inches(3.0), Inches(1.5), "Guest scope", [
        "Browse academies, trainers, and courses without registration.",
        "Read reviews and filter results before sign-up.",
    ], BLUE, "G")
    add_role_card(slide, Inches(7.22), Inches(1.6), Inches(3.0), Inches(1.5), "Registered scope", [
        "Book and cancel courses online.",
        "Track bookings and submit ratings.",
    ], GOOD, "U")
    add_role_card(slide, Inches(10.41), Inches(1.6), Inches(2.06), Inches(1.5), "Academy/admin", [
        "Dashboards, course approval, and oversight.",
    ], GOLD, "A")

    add_panel(slide, Inches(0.84), Inches(3.42), Inches(11.66), Inches(2.55), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.08), Inches(3.66), Inches(3.0), Inches(0.25), "In-scope capabilities", font_size=16, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(4.02), Inches(10.7), Inches(1.65), [
        "• Search, filter, and sort academies and courses",
        "• View academy profiles, trainer profiles, reviews, and ratings",
        "• Manage academy courses and trainers through a dashboard",
        "• Enforce course approval before publication",
        "• Focus on a web-based experience suitable for graduation defense",
    ], font_size=13, color=MUTED)

    add_footer(slide, 4)
    add_transition(slide)


def add_gap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Related Works & Gap", "The comparison in the report shows a missing centralized solution")

    features = [
        ("Multi-academy marketplace", "Learn 2 Drive is the only platform in the comparison matrix that aggregates multiple academies.", True),
        ("Online booking + dashboards", "Several tools support booking, but they do not combine academy dashboards, approval, and review controls.", False),
        ("Search and preference filters", "Most alternatives are limited to either basic booking or theoretical learning.", False),
        ("Transparent quality layer", "The platform adds ratings, reviews, and a course approval workflow that others lack.", True),
    ]

    for i, (title, desc, emphasis) in enumerate(features):
        left = Inches(0.85 + (i % 2) * 5.95)
        top = Inches(1.55 + (i // 2) * 1.95)
        fill = PANEL_2 if emphasis else PANEL
        accent = CYAN if emphasis else BLUE
        add_panel(slide, left, top, Inches(5.45), Inches(1.6), fill=fill, line=accent)
        add_text(slide, left + Inches(0.2), top + Inches(0.12), Inches(3.9), Inches(0.24), title, font_size=15, bold=True)
        add_text(slide, left + Inches(0.2), top + Inches(0.48), Inches(5.0), Inches(0.62), desc, font_size=11.5, color=MUTED)
        add_chip(slide, left + Inches(4.45), top + Inches(0.14), Inches(0.82), "Gap" if emphasis else "Limit", fill=accent, font_size=9)

    add_text(slide, Inches(0.92), Inches(5.92), Inches(11.2), Inches(0.55), "Result: Learn 2 Drive combines centralized discovery, online booking, academy management, review/rating visibility, and admin approval in a single platform.", font_size=13, color=WHITE, bold=True)
    add_footer(slide, 5)
    add_transition(slide)


def add_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Features & Users", "Four actors, each with a tailored experience")

    add_role_card(slide, Inches(0.84), Inches(1.48), Inches(2.8), Inches(2.1), "Guests", [
        "Browse academies and courses",
        "Open academy and trainer pages",
        "Read ratings and reviews",
        "Search and filter without login",
    ], CYAN, "G")
    add_role_card(slide, Inches(3.85), Inches(1.48), Inches(2.8), Inches(2.1), "Registered users", [
        "Book and cancel courses",
        "View booking history",
        "Submit ratings and reviews",
        "Manage profile information",
    ], BLUE, "U")
    add_role_card(slide, Inches(6.86), Inches(1.48), Inches(2.8), Inches(2.1), "Academies", [
        "Manage dashboards, courses, trainers",
        "Review bookings and feedback",
        "Submit new courses for approval",
        "Monitor public profile visibility",
    ], GOOD, "A")
    add_role_card(slide, Inches(9.87), Inches(1.48), Inches(2.56), Inches(2.1), "Administrators", [
        "Approve or reject courses",
        "Manage users and academies",
        "Delete inappropriate reviews",
        "Oversee platform reports",
    ], GOLD, "D")

    add_panel(slide, Inches(0.84), Inches(4.05), Inches(11.6), Inches(1.85), fill=PANEL_2, line=CYAN)
    add_text(slide, Inches(1.05), Inches(4.28), Inches(2.8), Inches(0.25), "User journey", font_size=15, bold=True)

    steps = ["Guest browse", "Register", "Book", "Review", "Academy manage", "Admin approve"]
    x = 1.05
    for i, step in enumerate(steps):
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.72), Inches(1.55), Inches(0.38))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(BLUE if i % 2 == 0 else CYAN)
        chip.line.fill.background()
        tf = chip.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = rgb(BG if i % 2 else WHITE)
        x += 1.84
        if i < len(steps) - 1:
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x - 0.28), Inches(4.91), Inches(x - 0.08), Inches(4.91))
            arrow.line.color.rgb = rgb(CYAN)
            arrow.line.width = Pt(1.5)

    add_footer(slide, 6)
    add_transition(slide)


def add_business_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Business Plan Highlights", "The report frames the platform as a scalable marketplace with a clear launch strategy")

    add_metric_card(slide, Inches(0.86), Inches(1.55), Inches(2.45), Inches(1.15), "Market focus", "Cairo + Giza", "Year 1 target region")
    add_metric_card(slide, Inches(3.50), Inches(1.55), Inches(2.45), Inches(1.15), "Target base", "175 academies", "Conservative addressable market")
    add_metric_card(slide, Inches(6.14), Inches(1.55), Inches(2.45), Inches(1.15), "Revenue mix", "Subs + fees + ads", "Diversified model")
    add_metric_card(slide, Inches(8.78), Inches(1.55), Inches(2.45), Inches(1.15), "Go-to-market", "BD + SEO", "Outreach-led growth")

    add_panel(slide, Inches(0.86), Inches(3.02), Inches(5.9), Inches(3.05), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.08), Inches(3.23), Inches(2.8), Inches(0.25), "Value proposition", font_size=16, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(3.63), Inches(5.1), Inches(2.1), [
        "• One digital place to discover and compare academies",
        "• Guests can browse without registration, lowering adoption friction",
        "• Quality control is enforced through course approval",
        "• Academies gain a managed dashboard and structured promotion channel",
        "• Learners gain trust through reviews, ratings, and filters",
    ], font_size=12.5, color=MUTED)

    add_panel(slide, Inches(7.02), Inches(3.02), Inches(5.25), Inches(3.05), fill=PANEL, line=CYAN)
    add_text(slide, Inches(7.25), Inches(3.23), Inches(2.8), Inches(0.25), "Business model", font_size=16, bold=True)
    add_multiline_text(slide, Inches(7.25), Inches(3.63), Inches(4.7), Inches(2.05), [
        "• Subscriptions from paying academies",
        "• Booking commission from active reservations",
        "• Secondary advertising revenue stream",
        "• Phase-based hiring and sales outreach",
        "• Clear path to scale after launch",
    ], font_size=12.5, color=MUTED)

    add_footer(slide, 7)
    add_transition(slide)


def add_financial_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Financial Snapshot", "Key values from the report's startup capital and Year 1 projection")

    add_metric_card(slide, Inches(0.86), Inches(1.52), Inches(2.35), Inches(1.08), "Startup capital", "EGP 180,000", "Fully funded by founders and family loan")
    add_metric_card(slide, Inches(3.42), Inches(1.52), Inches(2.35), Inches(1.08), "Year 1 revenue", "EGP 1,039,750", "Subscriptions, commissions, ads")
    add_metric_card(slide, Inches(5.98), Inches(1.52), Inches(2.35), Inches(1.08), "Year 1 profit", "EGP 63,570", "Cumulative P&L break-even by Month 12")
    add_metric_card(slide, Inches(8.54), Inches(1.52), Inches(2.35), Inches(1.08), "Monthly break-even", "Month 8", "First profitable month")
    add_metric_card(slide, Inches(11.10), Inches(1.52), Inches(1.08), Inches(1.08), "", "45", "paying academies")

    add_panel(slide, Inches(0.86), Inches(2.85), Inches(7.0), Inches(3.65), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.08), Inches(3.05), Inches(2.8), Inches(0.25), "Cash balance trajectory", font_size=15, bold=True)
    chart_data = ChartData()
    chart_data.categories = ["Start", "Day 0", "M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8", "M9", "M10", "M11", "M12"]
    chart_data.add_series("Cash balance", (180000, 144000, 112000, 85000, 58000, 37669, 32362, 38521, 22286, 31531, 68510, 140661, 267025, 454590))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1.02), Inches(3.35), Inches(6.5), Inches(2.8), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.line.color.rgb = rgb(CYAN)

    add_panel(slide, Inches(8.12), Inches(2.85), Inches(4.05), Inches(3.65), fill=PANEL, line=CYAN)
    add_text(slide, Inches(8.35), Inches(3.05), Inches(2.6), Inches(0.25), "Profitability notes", font_size=15, bold=True)
    add_multiline_text(slide, Inches(8.35), Inches(3.42), Inches(3.55), Inches(2.6), [
        "• First monthly profit appears in Month 8",
        "• Year 1 cumulative break-even is reached in Month 12",
        "• Cash balance bottoms at EGP 22,286 in Month 7",
        "• Year 1 includes EGP 247,020 non-cash deferred salary",
        "• Funding is conservative and based on a phased launch",
    ], font_size=12.3, color=MUTED)

    add_footer(slide, 8)
    add_transition(slide)


def add_analysis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Analysis", "Use case and activity logic summarized as a defense-friendly flow")

    add_panel(slide, Inches(0.84), Inches(1.5), Inches(4.2), Inches(5.05), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.05), Inches(1.73), Inches(2.0), Inches(0.25), "Use case view", font_size=16, bold=True)
    roles = [("Guest", CYAN), ("Registered user", BLUE), ("Academy", GOOD), ("Admin", GOLD)]
    for i, (label, accent) in enumerate(roles):
        y = 2.08 + i * 0.88
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.05), Inches(y), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(accent)
        circle.line.fill.background()
        add_text(slide, Inches(1.1), Inches(y + 0.03), Inches(0.32), Inches(0.18), str(i + 1), font_size=11, color=BG, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.62), Inches(y - 0.01), Inches(2.0), Inches(0.28), label, font_size=13, bold=True)
        add_text(slide, Inches(1.62), Inches(y + 0.25), Inches(2.45), Inches(0.35), ["browse, filter, read", "book, cancel, rate", "manage content, trainers", "approve, monitor, moderate"][i], font_size=10.5, color=MUTED)

    add_panel(slide, Inches(5.32), Inches(1.5), Inches(6.82), Inches(5.05), fill=PANEL, line=CYAN)
    add_text(slide, Inches(5.55), Inches(1.73), Inches(2.4), Inches(0.25), "Activity flow", font_size=16, bold=True)
    nodes = [
        ("Browse", "Academies and courses", CYAN),
        ("Filter", "By location, price, rating, gender", BLUE),
        ("Book", "Choose trainer, date, and time", GOOD),
        ("Approve", "Admin reviews academy course request", GOLD),
        ("Review", "User rates the completed experience", RED),
    ]
    for i, (title, desc, accent) in enumerate(nodes):
        y = 2.1 + i * 0.79
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(y), Inches(5.8), Inches(0.58))
        node.fill.solid()
        node.fill.fore_color.rgb = rgb(PANEL_2 if i % 2 == 0 else PANEL)
        node.line.color.rgb = rgb(accent)
        node.line.width = Pt(1.0)
        add_chip(slide, Inches(5.78), Inches(y + 0.13), Inches(0.72), title[:3].upper(), fill=accent, font_size=9)
        add_text(slide, Inches(6.65), Inches(y + 0.06), Inches(2.2), Inches(0.18), title, font_size=12.5, bold=True)
        add_text(slide, Inches(7.9), Inches(y + 0.06), Inches(3.1), Inches(0.18), desc, font_size=10.4, color=MUTED)
        if i < len(nodes) - 1:
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.45), Inches(y + 0.6), Inches(8.45), Inches(y + 0.73))
            arrow.line.color.rgb = rgb(CYAN)
            arrow.line.width = Pt(1.5)

    add_footer(slide, 9)
    add_transition(slide)


def add_design_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Design", "A simplified ERD and database mapping for the main entities")

    add_panel(slide, Inches(0.84), Inches(1.45), Inches(12.1), Inches(5.1), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.06), Inches(1.7), Inches(3.0), Inches(0.25), "Core entities", font_size=16, bold=True)

    center_x = 6.8
    center_y = 3.6
    nodes = [
        ("User", center_x - 1.8, center_y - 1.55, CYAN),
        ("Academy", center_x + 0.8, center_y - 1.25, BLUE),
        ("Trainer", center_x - 1.7, center_y + 0.65, GOOD),
        ("Course", center_x + 0.95, center_y + 0.75, GOLD),
        ("Booking", center_x - 0.25, center_y - 0.05, RED),
    ]
    for label, x, y, accent in nodes:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.35), Inches(0.62))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(accent)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = rgb(BG)

    conns = [
        ((center_x - 0.45, center_y - 1.0), (center_x + 0.2, center_y - 0.1)),
        ((center_x + 1.45, center_y - 0.65), (center_x + 0.35, center_y + 0.05)),
        ((center_x - 0.35, center_y + 0.65), (center_x - 0.02, center_y + 0.3)),
        ((center_x + 1.45, center_y + 1.05), (center_x + 0.4, center_y + 0.35)),
    ]
    for (x1, y1), (x2, y2) in conns:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = rgb(CYAN)
        line.line.width = Pt(1.25)

    labels = [
        ("Many-to-many", "Academy ↔ Location / Course ↔ Trainer"),
        ("One-to-many", "Academy → Trainers / Courses"),
        ("Booking rules", "Unique confirmed slot per trainer"),
        ("Content model", "Generic ratings and reviews"),
    ]
    for i, (title, desc) in enumerate(labels):
        left = Inches(1.05 + (i % 2) * 5.95)
        top = Inches(1.95 + (i // 2) * 1.52)
        add_panel(slide, left, top, Inches(4.9), Inches(1.08), fill=PANEL, line=BLUE)
        add_text(slide, left + Inches(0.18), top + Inches(0.1), Inches(2.0), Inches(0.2), title, font_size=12.5, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.38), Inches(4.3), Inches(0.3), desc, font_size=11, color=MUTED)

    add_footer(slide, 10)
    add_transition(slide)


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Architecture & Tools", "Client-server architecture with Django REST, React, JWT, and SQLite")

    layers = [
        ("Frontend", "React.js + Tailwind CSS", CYAN),
        ("API Layer", "Django REST Framework + JWT", BLUE),
        ("Data Layer", "SQLite + media + SMTP", GOOD),
    ]
    for i, (title, desc, accent) in enumerate(layers):
        top = Inches(1.45 + i * 1.6)
        add_panel(slide, Inches(0.92), top, Inches(7.0), Inches(1.2), fill=PANEL_2 if i != 1 else PANEL, line=accent)
        add_chip(slide, Inches(1.15), top + Inches(0.2), Inches(1.0), title, fill=accent, font_size=10)
        add_text(slide, Inches(2.35), top + Inches(0.12), Inches(4.8), Inches(0.25), desc, font_size=13, bold=True)
        add_text(slide, Inches(2.35), top + Inches(0.45), Inches(4.8), Inches(0.25), [
            "User interface and responsive pages",
            "Protected endpoints and business rules",
            "Database tables, uploaded media, emails",
        ][i], font_size=10.5, color=MUTED)

    arrow1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.45), Inches(2.65), Inches(4.45), Inches(3.0))
    arrow1.line.color.rgb = rgb(CYAN)
    arrow1.line.width = Pt(2)
    arrow2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.45), Inches(4.25), Inches(4.45), Inches(4.6))
    arrow2.line.color.rgb = rgb(CYAN)
    arrow2.line.width = Pt(2)

    add_panel(slide, Inches(8.2), Inches(1.45), Inches(4.0), Inches(4.25), fill=PANEL, line=CYAN)
    add_text(slide, Inches(8.45), Inches(1.68), Inches(2.5), Inches(0.25), "Tools used", font_size=16, bold=True)
    tools = ["Django REST Framework", "React.js", "Tailwind CSS", "JWT Authentication", "SQLite", "Gmail SMTP", "Scrum Agile"]
    x = 8.45
    y = 2.12
    for i, tool in enumerate(tools):
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.68), Inches(0.38))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(BLUE if i % 2 == 0 else CYAN)
        chip.line.fill.background()
        tf = chip.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = tool
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(9.3)
        p.font.bold = True
        p.font.color.rgb = rgb(BG if i % 2 else WHITE)
        x += 1.8
        if x > 10.6:
            x = 8.45
            y += 0.52

    add_text(slide, Inches(8.45), Inches(3.95), Inches(3.45), Inches(0.25), "Method", font_size=16, bold=True)
    add_text(slide, Inches(8.45), Inches(4.23), Inches(3.4), Inches(0.5), "Scrum-based development with iterative delivery and continuous adaptation.", font_size=12, color=MUTED)

    add_footer(slide, 11)
    add_transition(slide)


def add_implementation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Implementation Highlights", "The report shows the system through live screens and code-level screenshots")

    cards = [
        ("Custom auth", "Role-based user model and JWT login flow"),
        ("Public browsing", "Academy and course discovery with filters"),
        ("Booking flow", "Trainer selection, date picking, slot control"),
        ("Approval system", "Admin review before course publication"),
    ]
    for i, (title, desc) in enumerate(cards):
        left = Inches(0.86 + (i % 2) * 6.08)
        top = Inches(1.55 + (i // 2) * 1.72)
        add_panel(slide, left, top, Inches(5.66), Inches(1.38), fill=PANEL_2, line=CYAN if i % 2 == 0 else BLUE)
        add_chip(slide, left + Inches(0.18), top + Inches(0.18), Inches(0.72), str(i + 1), fill=CYAN if i % 2 == 0 else BLUE, font_size=10)
        add_text(slide, left + Inches(1.05), top + Inches(0.13), Inches(2.4), Inches(0.22), title, font_size=14, bold=True)
        add_text(slide, left + Inches(1.05), top + Inches(0.45), Inches(4.2), Inches(0.35), desc, font_size=11.2, color=MUTED)

    add_panel(slide, Inches(0.86), Inches(4.98), Inches(11.35), Inches(1.0), fill=PANEL, line=GOOD)
    add_text(slide, Inches(1.08), Inches(5.18), Inches(2.5), Inches(0.25), "Implementation stack seen in the report", font_size=14, bold=True)
    add_multiline_text(slide, Inches(4.3), Inches(5.08), Inches(7.6), Inches(0.46), [
        "• Django models, serializers, signals, permissions, URLs, and views",
        "• React pages, routing, Axios token handling, and protected routes",
    ], font_size=10.8, color=MUTED)

    add_footer(slide, 12)
    add_transition(slide)


def add_challenges_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Key Technical Challenges", "The hardest parts were mostly around state, constraints, and workflow integrity")

    challenges = [
        ("Availability logic", "Trainer working days and session windows had to be checked before booking."),
        ("Slot uniqueness", "Confirmed bookings needed a unique time slot constraint per trainer."),
        ("Role protection", "Different users needed different routes, permissions, and UI access."),
        ("Workflow automation", "Approval emails and status updates had to happen consistently."),
    ]
    for i, (title, desc) in enumerate(challenges):
        left = Inches(0.84 + (i % 2) * 5.96)
        top = Inches(1.55 + (i // 2) * 1.95)
        add_panel(slide, left, top, Inches(5.45), Inches(1.6), fill=PANEL_2, line=CYAN if i in (0, 3) else BLUE)
        add_chip(slide, left + Inches(0.18), top + Inches(0.13), Inches(0.95), "Challenge", fill=RED if i == 0 else GOLD if i == 1 else BLUE if i == 2 else CYAN, font_size=9)
        add_text(slide, left + Inches(1.3), top + Inches(0.12), Inches(2.2), Inches(0.22), title, font_size=14, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.5), Inches(5.0), Inches(0.48), desc, font_size=11.2, color=MUTED)

    add_text(slide, Inches(0.92), Inches(5.9), Inches(11.3), Inches(0.35), "The report indicates these problems were solved in the final system through validation, constraints, and backend automation.", font_size=12.5, color=WHITE, bold=True)
    add_footer(slide, 13)
    add_transition(slide)


def add_testing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Testing Results", "All 47 documented test cases passed successfully")

    add_metric_card(slide, Inches(0.86), Inches(1.5), Inches(2.7), Inches(1.2), "Pass rate", "47/47", "100% successful", accent=GOOD)
    add_metric_card(slide, Inches(3.80), Inches(1.5), Inches(2.7), Inches(1.2), "Method", "Manual + API", "REST Client and UI checks", accent=BLUE)
    add_metric_card(slide, Inches(6.74), Inches(1.5), Inches(2.7), Inches(1.2), "Coverage", "Auth → Reviews", "Core features validated", accent=CYAN)
    add_metric_card(slide, Inches(9.68), Inches(1.5), Inches(2.1), Inches(1.2), "Outcome", "Stable", "No critical gaps", accent=GOLD)

    add_panel(slide, Inches(0.86), Inches(3.0), Inches(5.9), Inches(3.1), fill=PANEL_2, line=GOOD)
    add_text(slide, Inches(1.08), Inches(3.22), Inches(2.0), Inches(0.25), "Test coverage", font_size=15, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(3.62), Inches(5.1), Inches(2.0), [
        "• Authentication and account flows",
        "• Academy registration and approval",
        "• Public academy, trainer, and course APIs",
        "• Booking creation and availability checks",
        "• Academy dashboard endpoints",
        "• Review and rating submission",
    ], font_size=12.2, color=MUTED)

    add_panel(slide, Inches(7.05), Inches(3.0), Inches(4.73), Inches(3.1), fill=PANEL, line=CYAN)
    add_text(slide, Inches(7.28), Inches(3.22), Inches(2.0), Inches(0.25), "Validation summary", font_size=15, bold=True)
    add_multiline_text(slide, Inches(7.28), Inches(3.62), Inches(4.0), Inches(1.95), [
        "• Positive and negative scenarios were both tested",
        "• REST Client was used for direct API requests",
        "• Frontend flows were checked manually in the browser",
        "• Results confirmed correct status codes and messages",
    ], font_size=12.2, color=MUTED)

    add_footer(slide, 14)
    add_transition(slide)


def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Conclusion & Future Work", "The project is ready for defense, with a clear roadmap for the next phase")

    add_panel(slide, Inches(0.86), Inches(1.5), Inches(5.85), Inches(4.95), fill=PANEL_2, line=CYAN)
    add_text(slide, Inches(1.08), Inches(1.74), Inches(2.2), Inches(0.25), "Conclusion", font_size=16, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(2.10), Inches(5.0), Inches(3.85), [
        "• Solves the lack of a centralized driving academy marketplace",
        "• Delivers a full web system built with Django REST and React",
        "• Includes secure authentication, approvals, bookings, and reviews",
        "• Passed all 47 documented test cases",
        "• Shows a viable business path with Month 8 profitability and Year 1 break-even",
    ], font_size=12.6, color=MUTED)

    add_panel(slide, Inches(7.0), Inches(1.5), Inches(4.78), Inches(4.95), fill=PANEL, line=GOOD)
    add_text(slide, Inches(7.22), Inches(1.74), Inches(2.2), Inches(0.25), "Future work", font_size=16, bold=True)
    futures = [
        "Online payment integration",
        "iOS and Android mobile apps",
        "Trainer self-registration portal",
        "Advanced analytics dashboard",
        "In-app messaging and AI recommendations",
        "PostgreSQL production migration and multilingual support",
    ]
    add_multiline_text(slide, Inches(7.22), Inches(2.10), Inches(4.1), Inches(3.55), [f"• {x}" for x in futures], font_size=11.9, color=MUTED)

    add_footer(slide, 15)
    add_transition(slide)


def add_thanks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title_band=True)
    add_text(slide, Inches(0.92), Inches(2.55), Inches(5.5), Inches(0.8), "Thank You", font_size=34, bold=True)
    add_text(slide, Inches(0.96), Inches(3.35), Inches(6.0), Inches(0.6), "Questions and discussion", font_size=16, color=MUTED)
    add_panel(slide, Inches(7.55), Inches(2.0), Inches(4.6), Inches(2.5), fill=PANEL, line=CYAN)
    add_text(slide, Inches(7.82), Inches(2.32), Inches(3.95), Inches(0.25), "Learn 2 Drive", font_size=20, bold=True)
    add_multiline_text(slide, Inches(7.82), Inches(2.82), Inches(3.8), Inches(1.0), [
        "A digital platform for driving academy discovery and booking",
        "2025–2026 Graduation Project",
    ], font_size=12.5, color=MUTED)


def build_presentation():
    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_agenda_slide(prs)
    add_problem_slide(prs)
    add_scope_slide(prs)
    add_gap_slide(prs)
    add_features_slide(prs)
    add_business_plan_slide(prs)
    add_financial_slide(prs)
    add_analysis_slide(prs)
    add_design_slide(prs)
    add_architecture_slide(prs)
    add_implementation_slide(prs)
    add_challenges_slide(prs)
    add_testing_slide(prs)
    add_conclusion_slide(prs)
    add_thanks_slide(prs)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Saved presentation to {OUTPUT}")


def add_background(slide, title_band=False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    width = prs.slide_width
    height = prs.slide_height

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, width, Inches(0.24))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(CYAN)
    bar.line.fill.background()

    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.25), Inches(0.35), Inches(0.08), Inches(6.65))
    side.fill.solid()
    side.fill.fore_color.rgb = rgb(BLUE)
    side.line.fill.background()

    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.9), Inches(-0.4), Inches(3.0), Inches(3.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = rgb(BLUE)
    glow.fill.transparency = 0.82
    glow.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.6), Inches(5.3), Inches(2.4), Inches(2.4))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = rgb(CYAN)
    glow2.fill.transparency = 0.86
    glow2.line.fill.background()

    if title_band:
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.15), Inches(5.6), Inches(0.6))
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(PANEL_2)
        band.line.color.rgb = rgb(CYAN)
        band.line.width = Pt(1.5)


def add_footer(slide, index):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(7.0), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(CYAN)
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.62), Inches(7.04), Inches(8.0), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Learn 2 Drive | BIS Graduation Project"
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.color.rgb = rgb(SLATE)

    page = slide.shapes.add_textbox(Inches(12.15), Inches(7.02), Inches(0.6), Inches(0.25))
    tf = page.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(index)
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = rgb(CYAN)


def add_heading(slide, title, subtitle=None, top=0.52, size=24):
    box = slide.shapes.add_textbox(Inches(0.62), Inches(top), Inches(8.8), Inches(0.85))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = rgb(WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Aptos"
        p2.font.size = Pt(11)
        p2.font.color.rgb = rgb(MUTED)
        p2.space_before = Pt(4)


def add_panel(slide, left, top, width, height, fill=PANEL, line=BLUE, radius=True):
    shape = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(fill)
    box.line.color.rgb = rgb(line)
    box.line.width = Pt(1.25)
    return box


def add_text(slide, left, top, width, height, text, font_size=16, color=WHITE, bold=False, font="Aptos", align=PP_ALIGN.LEFT, margin=0.08):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(font_size)
    p.font.color.rgb = rgb(color)
    p.font.bold = bold
    p.alignment = align
    return box


def add_multiline_text(slide, left, top, width, height, lines, font_size=15, color=WHITE, bullet=False, leading=1.1):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb(color)
        p.line_spacing = Pt(font_size * leading)
        if bullet:
            p.level = 0
            p.bullet = True
    return box


def add_metric_card(slide, left, top, width, height, label, value, caption=None, accent=CYAN):
    card = add_panel(slide, left, top, width, height, fill=PANEL_2, line=accent)
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.25), label, font_size=10, color=MUTED)
    add_text(slide, left + Inches(0.18), top + Inches(0.36), width - Inches(0.36), Inches(0.4), value, font_size=22, color=WHITE, bold=True)
    if caption:
        add_text(slide, left + Inches(0.18), top + Inches(0.9), width - Inches(0.36), Inches(0.4), caption, font_size=10, color=SLATE)
    return card


def add_chip(slide, left, top, width, text, fill=BLUE, font_size=10):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.32))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(fill)
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = rgb(WHITE)
    return chip


def add_role_card(slide, left, top, width, height, title, body, accent, symbol):
    add_panel(slide, left, top, width, height, fill=PANEL_2, line=accent)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), top + Inches(0.18), Inches(0.52), Inches(0.52))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(accent)
    badge.line.fill.background()
    add_text(slide, left + Inches(0.18), top + Inches(0.20), Inches(0.52), Inches(0.28), symbol, font_size=18, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.84), top + Inches(0.16), width - Inches(1.0), Inches(0.25), title, font_size=15, bold=True)
    add_multiline_text(slide, left + Inches(0.18), top + Inches(0.68), width - Inches(0.36), height - Inches(0.8), body, font_size=11, color=MUTED)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title_band=True)

    add_chip(slide, Inches(0.85), Inches(0.48), Inches(1.5), "GRADUATION DEFENSE", fill=BLUE)
    add_text(slide, Inches(0.82), Inches(1.38), Inches(7.6), Inches(1.0), "Learn 2 Drive", font_size=34, bold=True)
    add_text(slide, Inches(0.84), Inches(2.28), Inches(7.8), Inches(0.9), "A digital platform for driving academy discovery and booking", font_size=16, color=MUTED)

    add_metric_card(slide, Inches(0.86), Inches(3.05), Inches(2.0), Inches(1.15), "Project Year", "2025–2026", "BIS Graduation Project")
    add_metric_card(slide, Inches(3.08), Inches(3.05), Inches(2.0), Inches(1.15), "Team Members", "1. 2. 3. 4. 5.", "Insert final names")
    add_metric_card(slide, Inches(5.30), Inches(3.05), Inches(2.2), Inches(1.15), "Supervisors", "IT + Business", "Dr. ___________________")

    add_panel(slide, Inches(8.25), Inches(1.18), Inches(4.15), Inches(4.8), fill=PANEL, line=CYAN)
    add_text(slide, Inches(8.52), Inches(1.42), Inches(3.5), Inches(0.35), "Why this project matters", font_size=16, bold=True)
    add_multiline_text(slide, Inches(8.55), Inches(1.84), Inches(3.45), Inches(2.15), [
        "• Centralizes academy discovery in one place",
        "• Replaces manual booking with online flows",
        "• Adds approval, ratings, and transparent comparison",
        "• Supports guest browsing and role-based dashboards",
    ], font_size=13, color=MUTED)

    roadmap = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.52), Inches(4.25), Inches(3.45), Inches(1.15))
    roadmap.fill.solid()
    roadmap.fill.fore_color.rgb = rgb(BLUE)
    roadmap.line.fill.background()
    add_text(slide, Inches(8.7), Inches(4.44), Inches(3.1), Inches(0.25), "Modern defense-ready system design", font_size=14, bold=True)
    add_text(slide, Inches(8.7), Inches(4.72), Inches(3.0), Inches(0.20), "Dark cyan palette | clean visuals | clear flow", font_size=10, color=WHITE)

    add_footer(slide, 1)
    add_transition(slide)


def add_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Agenda", "15 content slides covering the project from problem to future work")

    items = [
        ("01", "Title", "Project identity and defense context"),
        ("02", "Agenda", "Roadmap of the presentation"),
        ("03", "Business Problem", "Current market pain points"),
        ("04", "Project Idea & Scope", "What the platform solves"),
        ("05", "Related Works & Gap", "Why Learn 2 Drive is different"),
        ("06", "System Features & Users", "Actors and responsibilities"),
        ("07", "Business Plan Highlights", "Go-to-market and value model"),
        ("08", "Financial Snapshot", "Startup capital and Year 1 outlook"),
        ("09", "System Analysis", "Use case and activity flows"),
        ("10", "System Design", "ERD and database mapping"),
        ("11", "Architecture & Tools", "Technology stack and layers"),
        ("12", "Implementation Highlights", "Core screens and code parts"),
        ("13", "Key Technical Challenges", "Hard parts and solutions"),
        ("14", "Testing Results", "47/47 test cases passed"),
        ("15", "Conclusion & Future Work", "Summary and next steps"),
    ]

    for idx, (num, title, desc) in enumerate(items):
        row = idx // 3
        col = idx % 3
        left = Inches(0.82 + col * 4.12)
        top = Inches(1.35 + row * 2.0)
        add_panel(slide, left, top, Inches(3.76), Inches(1.55), fill=PANEL_2, line=BLUE)
        add_chip(slide, left + Inches(0.16), top + Inches(0.14), Inches(0.52), num, fill=CYAN, font_size=10)
        add_text(slide, left + Inches(0.84), top + Inches(0.12), Inches(2.6), Inches(0.25), title, font_size=15, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.60), Inches(3.25), Inches(0.55), desc, font_size=11, color=MUTED)

    add_footer(slide, 2)
    add_transition(slide)


def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Business Problem", "The report identifies a fragmented, manual market with limited transparency")

    add_panel(slide, Inches(0.82), Inches(1.52), Inches(6.05), Inches(4.9), fill=PANEL_2, line=CYAN)
    add_text(slide, Inches(1.05), Inches(1.72), Inches(2.2), Inches(0.25), "Core pain points", font_size=16, bold=True)
    bullets = [
        "Information is scattered across websites, social media, and word-of-mouth.",
        "Pricing, packages, and instructor quality are often not transparent.",
        "Booking is still manual and usually requires phone calls or in-person visits.",
        "Learners with specific needs, such as female instructors, cannot filter easily.",
        "There is no reliable review layer for comparing academy quality.",
    ]
    add_multiline_text(slide, Inches(1.05), Inches(2.08), Inches(5.45), Inches(3.55), [f"• {x}" for x in bullets], font_size=13, color=MUTED)

    add_panel(slide, Inches(7.1), Inches(1.52), Inches(5.35), Inches(4.9), fill=PANEL, line=BLUE)
    add_text(slide, Inches(7.35), Inches(1.72), Inches(2.6), Inches(0.25), "Business impact", font_size=16, bold=True)
    add_metric_card(slide, Inches(7.35), Inches(2.05), Inches(2.25), Inches(1.0), "Transparency", "Low", "Hidden pricing and service quality", accent=RED)
    add_metric_card(slide, Inches(9.75), Inches(2.05), Inches(2.25), Inches(1.0), "Booking", "Manual", "Offline, time-consuming process", accent=GOLD)
    add_metric_card(slide, Inches(7.35), Inches(3.20), Inches(2.25), Inches(1.0), "Search", "Fragmented", "No central comparison layer", accent=BLUE)
    add_metric_card(slide, Inches(9.75), Inches(3.20), Inches(2.25), Inches(1.0), "Feedback", "Missing", "No trust-building review system", accent=CYAN)

    add_text(slide, Inches(7.35), Inches(4.55), Inches(4.55), Inches(0.55), "The platform answers this gap with a centralized marketplace, online booking, filtering, ratings, and approval control.", font_size=13, color=MUTED)
    add_footer(slide, 3)
    add_transition(slide)


def add_scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Project Idea & Scope", "A centralized platform connecting learners, academies, trainers, and administrators")

    add_role_card(slide, Inches(0.84), Inches(1.6), Inches(3.0), Inches(1.5), "Central platform", [
        "Connects driving learners with academies in one online environment.",
        "Supports discovery, comparison, and booking.",
    ], CYAN, "C")
    add_role_card(slide, Inches(4.03), Inches(1.6), Inches(3.0), Inches(1.5), "Guest scope", [
        "Browse academies, trainers, and courses without registration.",
        "Read reviews and filter results before sign-up.",
    ], BLUE, "G")
    add_role_card(slide, Inches(7.22), Inches(1.6), Inches(3.0), Inches(1.5), "Registered scope", [
        "Book and cancel courses online.",
        "Track bookings and submit ratings.",
    ], GOOD, "U")
    add_role_card(slide, Inches(10.41), Inches(1.6), Inches(2.06), Inches(1.5), "Academy/admin", [
        "Dashboards, course approval, and oversight.",
    ], GOLD, "A")

    add_panel(slide, Inches(0.84), Inches(3.42), Inches(11.66), Inches(2.55), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.08), Inches(3.66), Inches(3.0), Inches(0.25), "In-scope capabilities", font_size=16, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(4.02), Inches(10.7), Inches(1.65), [
        "• Search, filter, and sort academies and courses",
        "• View academy profiles, trainer profiles, reviews, and ratings",
        "• Manage academy courses and trainers through a dashboard",
        "• Enforce course approval before publication",
        "• Focus on a web-based experience suitable for graduation defense",
    ], font_size=13, color=MUTED)

    add_footer(slide, 4)
    add_transition(slide)


def add_gap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Related Works & Gap", "The comparison in the report shows a missing centralized solution")

    features = [
        ("Multi-academy marketplace", "Learn 2 Drive is the only platform in the comparison matrix that aggregates multiple academies.", True),
        ("Online booking + dashboards", "Several tools support booking, but they do not combine academy dashboards, approval, and review controls.", False),
        ("Search and preference filters", "Most alternatives are limited to either basic booking or theoretical learning.", False),
        ("Transparent quality layer", "The platform adds ratings, reviews, and a course approval workflow that others lack.", True),
    ]

    for i, (title, desc, emphasis) in enumerate(features):
        left = Inches(0.85 + (i % 2) * 5.95)
        top = Inches(1.55 + (i // 2) * 1.95)
        fill = PANEL_2 if emphasis else PANEL
        accent = CYAN if emphasis else BLUE
        add_panel(slide, left, top, Inches(5.45), Inches(1.6), fill=fill, line=accent)
        add_text(slide, left + Inches(0.2), top + Inches(0.12), Inches(3.9), Inches(0.24), title, font_size=15, bold=True)
        add_text(slide, left + Inches(0.2), top + Inches(0.48), Inches(5.0), Inches(0.62), desc, font_size=11.5, color=MUTED)
        add_chip(slide, left + Inches(4.45), top + Inches(0.14), Inches(0.82), "Gap" if emphasis else "Limit", fill=accent, font_size=9)

    add_text(slide, Inches(0.92), Inches(5.92), Inches(11.2), Inches(0.55), "Result: Learn 2 Drive combines centralized discovery, online booking, academy management, review/rating visibility, and admin approval in a single platform.", font_size=13, color=WHITE, bold=True)
    add_footer(slide, 5)
    add_transition(slide)


def add_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Features & Users", "Four actors, each with a tailored experience")

    add_role_card(slide, Inches(0.84), Inches(1.48), Inches(2.8), Inches(2.1), "Guests", [
        "Browse academies and courses",
        "Open academy and trainer pages",
        "Read ratings and reviews",
        "Search and filter without login",
    ], CYAN, "G")
    add_role_card(slide, Inches(3.85), Inches(1.48), Inches(2.8), Inches(2.1), "Registered users", [
        "Book and cancel courses",
        "View booking history",
        "Submit ratings and reviews",
        "Manage profile information",
    ], BLUE, "U")
    add_role_card(slide, Inches(6.86), Inches(1.48), Inches(2.8), Inches(2.1), "Academies", [
        "Manage dashboards, courses, trainers",
        "Review bookings and feedback",
        "Submit new courses for approval",
        "Monitor public profile visibility",
    ], GOOD, "A")
    add_role_card(slide, Inches(9.87), Inches(1.48), Inches(2.56), Inches(2.1), "Administrators", [
        "Approve or reject courses",
        "Manage users and academies",
        "Delete inappropriate reviews",
        "Oversee platform reports",
    ], GOLD, "D")

    add_panel(slide, Inches(0.84), Inches(4.05), Inches(11.6), Inches(1.85), fill=PANEL_2, line=CYAN)
    add_text(slide, Inches(1.05), Inches(4.28), Inches(2.8), Inches(0.25), "User journey", font_size=15, bold=True)

    steps = ["Guest browse", "Register", "Book", "Review", "Academy manage", "Admin approve"]
    x = 1.05
    for i, step in enumerate(steps):
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.72), Inches(1.55), Inches(0.38))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(BLUE if i % 2 == 0 else CYAN)
        chip.line.fill.background()
        tf = chip.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = rgb(BG if i % 2 else WHITE)
        x += 1.84
        if i < len(steps) - 1:
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x - 0.28), Inches(4.91), Inches(x - 0.08), Inches(4.91))
            arrow.line.color.rgb = rgb(CYAN)
            arrow.line.width = Pt(1.5)

    add_footer(slide, 6)
    add_transition(slide)


def add_business_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Business Plan Highlights", "The report frames the platform as a scalable marketplace with a clear launch strategy")

    add_metric_card(slide, Inches(0.86), Inches(1.55), Inches(2.45), Inches(1.15), "Market focus", "Cairo + Giza", "Year 1 target region")
    add_metric_card(slide, Inches(3.50), Inches(1.55), Inches(2.45), Inches(1.15), "Target base", "175 academies", "Conservative addressable market")
    add_metric_card(slide, Inches(6.14), Inches(1.55), Inches(2.45), Inches(1.15), "Revenue mix", "Subs + fees + ads", "Diversified model")
    add_metric_card(slide, Inches(8.78), Inches(1.55), Inches(2.45), Inches(1.15), "Go-to-market", "BD + SEO", "Outreach-led growth")

    add_panel(slide, Inches(0.86), Inches(3.02), Inches(5.9), Inches(3.05), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.08), Inches(3.23), Inches(2.8), Inches(0.25), "Value proposition", font_size=16, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(3.63), Inches(5.1), Inches(2.1), [
        "• One digital place to discover and compare academies",
        "• Guests can browse without registration, lowering adoption friction",
        "• Quality control is enforced through course approval",
        "• Academies gain a managed dashboard and structured promotion channel",
        "• Learners gain trust through reviews, ratings, and filters",
    ], font_size=12.5, color=MUTED)

    add_panel(slide, Inches(7.02), Inches(3.02), Inches(5.25), Inches(3.05), fill=PANEL, line=CYAN)
    add_text(slide, Inches(7.25), Inches(3.23), Inches(2.8), Inches(0.25), "Business model", font_size=16, bold=True)
    add_multiline_text(slide, Inches(7.25), Inches(3.63), Inches(4.7), Inches(2.05), [
        "• Subscriptions from paying academies",
        "• Booking commission from active reservations",
        "• Secondary advertising revenue stream",
        "• Phase-based hiring and sales outreach",
        "• Clear path to scale after launch",
    ], font_size=12.5, color=MUTED)

    add_footer(slide, 7)
    add_transition(slide)


def add_financial_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Financial Snapshot", "Key values from the report's startup capital and Year 1 projection")

    add_metric_card(slide, Inches(0.86), Inches(1.52), Inches(2.35), Inches(1.08), "Startup capital", "EGP 180,000", "Fully funded by founders and family loan", accent=GOOD)
    add_metric_card(slide, Inches(3.42), Inches(1.52), Inches(2.35), Inches(1.08), "Year 1 revenue", "EGP 1,039,750", "Subscriptions, commissions, ads", accent=BLUE)
    add_metric_card(slide, Inches(5.98), Inches(1.52), Inches(2.35), Inches(1.08), "Year 1 profit", "EGP 63,570", "Cumulative P&L break-even by Month 12", accent=CYAN)
    add_metric_card(slide, Inches(8.54), Inches(1.52), Inches(2.35), Inches(1.08), "Monthly break-even", "Month 8", "First profitable month", accent=GOLD)
    add_metric_card(slide, Inches(11.10), Inches(1.52), Inches(1.08), Inches(1.08), "", "45", "paying academies", accent=RED)

    add_panel(slide, Inches(0.86), Inches(2.85), Inches(7.0), Inches(3.65), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.08), Inches(3.05), Inches(2.8), Inches(0.25), "Cash balance trajectory", font_size=15, bold=True)
    chart_data = ChartData()
    chart_data.categories = ["Start", "Day 0", "M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8", "M9", "M10", "M11", "M12"]
    chart_data.add_series("Cash balance", (180000, 144000, 112000, 85000, 58000, 37669, 32362, 38521, 22286, 31531, 68510, 140661, 267025, 454590))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(1.02), Inches(3.35), Inches(6.5), Inches(2.8), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.line.color.rgb = rgb(CYAN)

    add_panel(slide, Inches(8.12), Inches(2.85), Inches(4.05), Inches(3.65), fill=PANEL, line=CYAN)
    add_text(slide, Inches(8.35), Inches(3.05), Inches(2.6), Inches(0.25), "Profitability notes", font_size=15, bold=True)
    add_multiline_text(slide, Inches(8.35), Inches(3.42), Inches(3.55), Inches(2.6), [
        "• First monthly profit appears in Month 8",
        "• Year 1 cumulative break-even is reached in Month 12",
        "• Cash balance bottoms at EGP 22,286 in Month 7",
        "• Year 1 includes EGP 247,020 non-cash deferred salary",
        "• Funding is conservative and based on a phased launch",
    ], font_size=12.3, color=MUTED)

    add_footer(slide, 8)
    add_transition(slide)


def add_analysis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Analysis", "Use case and activity logic summarized as a defense-friendly flow")

    add_panel(slide, Inches(0.84), Inches(1.5), Inches(4.2), Inches(5.05), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.05), Inches(1.73), Inches(2.0), Inches(0.25), "Use case view", font_size=16, bold=True)
    roles = [("Guest", CYAN), ("Registered user", BLUE), ("Academy", GOOD), ("Admin", GOLD)]
    for i, (label, accent) in enumerate(roles):
        y = 2.08 + i * 0.88
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.05), Inches(y), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(accent)
        circle.line.fill.background()
        add_text(slide, Inches(1.1), Inches(y + 0.03), Inches(0.32), Inches(0.18), str(i + 1), font_size=11, color=BG, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.62), Inches(y - 0.01), Inches(2.0), Inches(0.28), label, font_size=13, bold=True)
        add_text(slide, Inches(1.62), Inches(y + 0.25), Inches(2.45), Inches(0.35), ["browse, filter, read", "book, cancel, rate", "manage content, trainers", "approve, monitor, moderate"][i], font_size=10.5, color=MUTED)

    add_panel(slide, Inches(5.32), Inches(1.5), Inches(6.82), Inches(5.05), fill=PANEL, line=CYAN)
    add_text(slide, Inches(5.55), Inches(1.73), Inches(2.4), Inches(0.25), "Activity flow", font_size=16, bold=True)
    nodes = [
        ("Browse", "Academies and courses", CYAN),
        ("Filter", "By location, price, rating, gender", BLUE),
        ("Book", "Choose trainer, date, and time", GOOD),
        ("Approve", "Admin reviews academy course request", GOLD),
        ("Review", "User rates the completed experience", RED),
    ]
    for i, (title, desc, accent) in enumerate(nodes):
        y = 2.1 + i * 0.79
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(y), Inches(5.8), Inches(0.58))
        node.fill.solid()
        node.fill.fore_color.rgb = rgb(PANEL_2 if i % 2 == 0 else PANEL)
        node.line.color.rgb = rgb(accent)
        node.line.width = Pt(1.0)
        add_chip(slide, Inches(5.78), Inches(y + 0.13), Inches(0.72), title[:3].upper(), fill=accent, font_size=9)
        add_text(slide, Inches(6.65), Inches(y + 0.06), Inches(2.2), Inches(0.18), title, font_size=12.5, bold=True)
        add_text(slide, Inches(7.9), Inches(y + 0.06), Inches(3.1), Inches(0.18), desc, font_size=10.4, color=MUTED)
        if i < len(nodes) - 1:
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.45), Inches(y + 0.6), Inches(8.45), Inches(y + 0.73))
            arrow.line.color.rgb = rgb(CYAN)
            arrow.line.width = Pt(1.5)

    add_footer(slide, 9)
    add_transition(slide)


def add_design_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Design", "A simplified ERD and database mapping for the main entities")

    add_panel(slide, Inches(0.84), Inches(1.45), Inches(12.1), Inches(5.1), fill=PANEL_2, line=BLUE)
    add_text(slide, Inches(1.06), Inches(1.7), Inches(3.0), Inches(0.25), "Core entities", font_size=16, bold=True)

    center_x = 6.8
    center_y = 3.6
    nodes = [
        ("User", center_x - 1.8, center_y - 1.55, CYAN),
        ("Academy", center_x + 0.8, center_y - 1.25, BLUE),
        ("Trainer", center_x - 1.7, center_y + 0.65, GOOD),
        ("Course", center_x + 0.95, center_y + 0.75, GOLD),
        ("Booking", center_x - 0.25, center_y - 0.05, RED),
    ]
    for label, x, y, accent in nodes:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.35), Inches(0.62))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(accent)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = rgb(BG)

    conns = [
        ((center_x - 0.45, center_y - 1.0), (center_x + 0.2, center_y - 0.1)),
        ((center_x + 1.45, center_y - 0.65), (center_x + 0.35, center_y + 0.05)),
        ((center_x - 0.35, center_y + 0.65), (center_x - 0.02, center_y + 0.3)),
        ((center_x + 1.45, center_y + 1.05), (center_x + 0.4, center_y + 0.35)),
    ]
    for (x1, y1), (x2, y2) in conns:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = rgb(CYAN)
        line.line.width = Pt(1.25)

    labels = [
        ("Many-to-many", "Academy ↔ Location / Course ↔ Trainer"),
        ("One-to-many", "Academy → Trainers / Courses"),
        ("Booking rules", "Unique confirmed slot per trainer"),
        ("Content model", "Generic ratings and reviews"),
    ]
    for i, (title, desc) in enumerate(labels):
        left = Inches(1.05 + (i % 2) * 5.95)
        top = Inches(1.95 + (i // 2) * 1.52)
        add_panel(slide, left, top, Inches(4.9), Inches(1.08), fill=PANEL, line=BLUE)
        add_text(slide, left + Inches(0.18), top + Inches(0.1), Inches(2.0), Inches(0.2), title, font_size=12.5, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.38), Inches(4.3), Inches(0.3), desc, font_size=11, color=MUTED)

    add_footer(slide, 10)
    add_transition(slide)


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "System Architecture & Tools", "Client-server architecture with Django REST, React, JWT, and SQLite")

    layers = [
        ("Frontend", "React.js + Tailwind CSS", CYAN),
        ("API Layer", "Django REST Framework + JWT", BLUE),
        ("Data Layer", "SQLite + media + SMTP", GOOD),
    ]
    for i, (title, desc, accent) in enumerate(layers):
        top = Inches(1.45 + i * 1.6)
        add_panel(slide, Inches(0.92), top, Inches(7.0), Inches(1.2), fill=PANEL_2 if i != 1 else PANEL, line=accent)
        add_chip(slide, Inches(1.15), top + Inches(0.2), Inches(1.0), title, fill=accent, font_size=10)
        add_text(slide, Inches(2.35), top + Inches(0.12), Inches(4.8), Inches(0.25), desc, font_size=13, bold=True)
        add_text(slide, Inches(2.35), top + Inches(0.45), Inches(4.8), Inches(0.25), [
            "User interface and responsive pages",
            "Protected endpoints and business rules",
            "Database tables, uploaded media, emails",
        ][i], font_size=10.5, color=MUTED)

    arrow1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.45), Inches(2.65), Inches(4.45), Inches(3.0))
    arrow1.line.color.rgb = rgb(CYAN)
    arrow1.line.width = Pt(2)
    arrow2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.45), Inches(4.25), Inches(4.45), Inches(4.6))
    arrow2.line.color.rgb = rgb(CYAN)
    arrow2.line.width = Pt(2)

    add_panel(slide, Inches(8.2), Inches(1.45), Inches(4.0), Inches(4.25), fill=PANEL, line=CYAN)
    add_text(slide, Inches(8.45), Inches(1.68), Inches(2.5), Inches(0.25), "Tools used", font_size=16, bold=True)
    tools = ["Django REST Framework", "React.js", "Tailwind CSS", "JWT Authentication", "SQLite", "Gmail SMTP", "Scrum Agile"]
    x = 8.45
    y = 2.12
    for i, tool in enumerate(tools):
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.68), Inches(0.38))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(BLUE if i % 2 == 0 else CYAN)
        chip.line.fill.background()
        tf = chip.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = tool
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(9.3)
        p.font.bold = True
        p.font.color.rgb = rgb(BG if i % 2 else WHITE)
        x += 1.8
        if x > 10.6:
            x = 8.45
            y += 0.52

    add_text(slide, Inches(8.45), Inches(3.95), Inches(3.45), Inches(0.25), "Method", font_size=16, bold=True)
    add_text(slide, Inches(8.45), Inches(4.23), Inches(3.4), Inches(0.5), "Scrum-based development with iterative delivery and continuous adaptation.", font_size=12, color=MUTED)

    add_footer(slide, 11)
    add_transition(slide)


def add_implementation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Implementation Highlights", "The report shows the system through live screens and code-level screenshots")

    cards = [
        ("Custom auth", "Role-based user model and JWT login flow"),
        ("Public browsing", "Academy and course discovery with filters"),
        ("Booking flow", "Trainer selection, date picking, slot control"),
        ("Approval system", "Admin review before course publication"),
    ]
    for i, (title, desc) in enumerate(cards):
        left = Inches(0.86 + (i % 2) * 6.08)
        top = Inches(1.55 + (i // 2) * 1.72)
        add_panel(slide, left, top, Inches(5.66), Inches(1.38), fill=PANEL_2, line=CYAN if i % 2 == 0 else BLUE)
        add_chip(slide, left + Inches(0.18), top + Inches(0.18), Inches(0.72), str(i + 1), fill=CYAN if i % 2 == 0 else BLUE, font_size=10)
        add_text(slide, left + Inches(1.05), top + Inches(0.13), Inches(2.4), Inches(0.22), title, font_size=14, bold=True)
        add_text(slide, left + Inches(1.05), top + Inches(0.45), Inches(4.2), Inches(0.35), desc, font_size=11.2, color=MUTED)

    add_panel(slide, Inches(0.86), Inches(4.98), Inches(11.35), Inches(1.0), fill=PANEL, line=GOOD)
    add_text(slide, Inches(1.08), Inches(5.18), Inches(2.5), Inches(0.25), "Implementation stack seen in the report", font_size=14, bold=True)
    add_multiline_text(slide, Inches(4.3), Inches(5.08), Inches(7.6), Inches(0.46), [
        "• Django models, serializers, signals, permissions, URLs, and views",
        "• React pages, routing, Axios token handling, and protected routes",
    ], font_size=10.8, color=MUTED)

    add_footer(slide, 12)
    add_transition(slide)


def add_challenges_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Key Technical Challenges", "The hardest parts were mostly around state, constraints, and workflow integrity")

    challenges = [
        ("Availability logic", "Trainer working days and session windows had to be checked before booking."),
        ("Slot uniqueness", "Confirmed bookings needed a unique time slot constraint per trainer."),
        ("Role protection", "Different users needed different routes, permissions, and UI access."),
        ("Workflow automation", "Approval emails and status updates had to happen consistently."),
    ]
    for i, (title, desc) in enumerate(challenges):
        left = Inches(0.84 + (i % 2) * 5.96)
        top = Inches(1.55 + (i // 2) * 1.95)
        add_panel(slide, left, top, Inches(5.45), Inches(1.6), fill=PANEL_2, line=CYAN if i in (0, 3) else BLUE)
        add_chip(slide, left + Inches(0.18), top + Inches(0.13), Inches(0.95), "Challenge", fill=RED if i == 0 else GOLD if i == 1 else BLUE if i == 2 else CYAN, font_size=9)
        add_text(slide, left + Inches(1.3), top + Inches(0.12), Inches(2.2), Inches(0.22), title, font_size=14, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.5), Inches(5.0), Inches(0.48), desc, font_size=11.2, color=MUTED)

    add_text(slide, Inches(0.92), Inches(5.9), Inches(11.3), Inches(0.35), "The report indicates these problems were solved in the final system through validation, constraints, and backend automation.", font_size=12.5, color=WHITE, bold=True)
    add_footer(slide, 13)
    add_transition(slide)


def add_testing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Testing Results", "All 47 documented test cases passed successfully")

    add_metric_card(slide, Inches(0.86), Inches(1.5), Inches(2.7), Inches(1.2), "Pass rate", "47/47", "100% successful", accent=GOOD)
    add_metric_card(slide, Inches(3.80), Inches(1.5), Inches(2.7), Inches(1.2), "Method", "Manual + API", "REST Client and UI checks", accent=BLUE)
    add_metric_card(slide, Inches(6.74), Inches(1.5), Inches(2.7), Inches(1.2), "Coverage", "Auth → Reviews", "Core features validated", accent=CYAN)
    add_metric_card(slide, Inches(9.68), Inches(1.5), Inches(2.1), Inches(1.2), "Outcome", "Stable", "No critical gaps", accent=GOLD)

    add_panel(slide, Inches(0.86), Inches(3.0), Inches(5.9), Inches(3.1), fill=PANEL_2, line=GOOD)
    add_text(slide, Inches(1.08), Inches(3.22), Inches(2.0), Inches(0.25), "Test coverage", font_size=15, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(3.62), Inches(5.1), Inches(2.0), [
        "• Authentication and account flows",
        "• Academy registration and approval",
        "• Public academy, trainer, and course APIs",
        "• Booking creation and availability checks",
        "• Academy dashboard endpoints",
        "• Review and rating submission",
    ], font_size=12.2, color=MUTED)

    add_panel(slide, Inches(7.05), Inches(3.0), Inches(4.73), Inches(3.1), fill=PANEL, line=CYAN)
    add_text(slide, Inches(7.28), Inches(3.22), Inches(2.0), Inches(0.25), "Validation summary", font_size=15, bold=True)
    add_multiline_text(slide, Inches(7.28), Inches(3.62), Inches(4.0), Inches(1.95), [
        "• Positive and negative scenarios were both tested",
        "• REST Client was used for direct API requests",
        "• Frontend flows were checked manually in the browser",
        "• Results confirmed correct status codes and messages",
    ], font_size=12.2, color=MUTED)

    add_footer(slide, 14)
    add_transition(slide)


def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_heading(slide, "Conclusion & Future Work", "The project is ready for defense, with a clear roadmap for the next phase")

    add_panel(slide, Inches(0.86), Inches(1.5), Inches(5.85), Inches(4.95), fill=PANEL_2, line=CYAN)
    add_text(slide, Inches(1.08), Inches(1.74), Inches(2.2), Inches(0.25), "Conclusion", font_size=16, bold=True)
    add_multiline_text(slide, Inches(1.08), Inches(2.10), Inches(5.0), Inches(3.85), [
        "• Solves the lack of a centralized driving academy marketplace",
        "• Delivers a full web system built with Django REST and React",
        "• Includes secure authentication, approvals, bookings, and reviews",
        "• Passed all 47 documented test cases",
        "• Shows a viable business path with Month 8 profitability and Year 1 break-even",
    ], font_size=12.6, color=MUTED)

    add_panel(slide, Inches(7.0), Inches(1.5), Inches(4.78), Inches(4.95), fill=PANEL, line=GOOD)
    add_text(slide, Inches(7.22), Inches(1.74), Inches(2.2), Inches(0.25), "Future work", font_size=16, bold=True)
    futures = [
        "Online payment integration",
        "iOS and Android mobile apps",
        "Trainer self-registration portal",
        "Advanced analytics dashboard",
        "In-app messaging and AI recommendations",
        "PostgreSQL production migration and multilingual support",
    ]
    add_multiline_text(slide, Inches(7.22), Inches(2.10), Inches(4.1), Inches(3.55), [f"• {x}" for x in futures], font_size=11.9, color=MUTED)

    add_footer(slide, 15)
    add_transition(slide)


def add_thanks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title_band=True)
    add_text(slide, Inches(0.92), Inches(2.55), Inches(5.5), Inches(0.8), "Thank You", font_size=34, bold=True)
    add_text(slide, Inches(0.96), Inches(3.35), Inches(6.0), Inches(0.6), "Questions and discussion", font_size=16, color=MUTED)
    add_panel(slide, Inches(7.55), Inches(2.0), Inches(4.6), Inches(2.5), fill=PANEL, line=CYAN)
    add_text(slide, Inches(7.82), Inches(2.32), Inches(3.95), Inches(0.25), "Learn 2 Drive", font_size=20, bold=True)
    add_multiline_text(slide, Inches(7.82), Inches(2.82), Inches(3.8), Inches(1.0), [
        "A digital platform for driving academy discovery and booking",
        "2025–2026 Graduation Project",
    ], font_size=12.5, color=MUTED)


def build_presentation():
    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_agenda_slide(prs)
    add_problem_slide(prs)
    add_scope_slide(prs)
    add_gap_slide(prs)
    add_features_slide(prs)
    add_business_plan_slide(prs)
    add_financial_slide(prs)
    add_analysis_slide(prs)
    add_design_slide(prs)
    add_architecture_slide(prs)
    add_implementation_slide(prs)
    add_challenges_slide(prs)
    add_testing_slide(prs)
    add_conclusion_slide(prs)
    add_thanks_slide(prs)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Saved presentation to {OUTPUT}")